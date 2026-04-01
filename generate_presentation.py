"""
Generates a professional 6-slide PowerPoint presentation for the
Global GHG Emissions Intelligence System capstone project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x02, 0x06, 0x17)   # #020617 — midnight navy
BG_CARD       = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a — slate-900
TEAL          = RGBColor(0x14, 0xB8, 0xA6)   # #14b8a6
EMERALD       = RGBColor(0x10, 0xB9, 0x81)   # #10b981
AMBER         = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
SLATE_300     = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_400     = RGBColor(0x94, 0xA3, 0xB8)
ACCENT_VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
RED           = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────
def add_filled_rect(slide, left, top, width, height, color, alpha_pct=100):
    """Add a solid filled rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    # Transparency
    if alpha_pct < 100:
        sp = shape._element
        solidFill = sp.spPr.find(qn('a:solidFill'))
        if solidFill is None:
            solidFill = etree.SubElement(sp.spPr, qn('a:solidFill'))
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is None:
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', f'{color[0]:02X}{color[1]:02X}{color[2]:02X}')
        alpha_elem = srgbClr.find(qn('a:alpha'))
        if alpha_elem is None:
            alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(alpha_pct * 1000)))
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label_value(slide, label, value, left, top,
                    label_color=SLATE_400, value_color=TEAL,
                    label_size=9, value_size=22):
    """Add a stacked label + value block."""
    add_text(slide, label.upper(), left, top, Inches(2.5), Inches(0.3),
             font_size=label_size, color=label_color, bold=True)
    add_text(slide, value, left, top + Inches(0.3), Inches(2.5), Inches(0.45),
             font_size=value_size, bold=True, color=value_color)


def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """Simulate gradient by stacking two overlapping rects."""
    add_filled_rect(slide, left, top, width, height, color1)
    add_filled_rect(slide, left + width // 2, top, width // 2, height, color2)


def add_pill_tag(slide, text, left, top, bg_color, text_color):
    """Add a small rounded-rectangle pill / badge."""
    shape = slide.shapes.add_shape(
        5,  # rounded rectangle (MSO_SHAPE.ROUNDED_RECTANGLE)
        left, top, Inches(1.8), Inches(0.32)
    )
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.color.rgb = text_color
    run.font.name = "Calibri"
    return shape


def add_section_line(slide, left, top, width, color=TEAL):
    bar = slide.shapes.add_shape(1, left, top, width, Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ─── Slide 1 — Title / Hero ───────────────────────────────────────────────────
def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, BG_DARK)

    # Left accent bar
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)

    # Top-right dark card block
    add_filled_rect(slide, Inches(8.5), 0, Inches(4.83), SLIDE_H, BG_CARD)

    # Badge
    add_pill_tag(slide, "CAPSTONE PROJECT", Inches(0.5), Inches(1.6),
                 RGBColor(0x0D, 0x3D, 0x38), TEAL)

    # Title
    add_text(slide, "Global GHG Emissions",
             Inches(0.5), Inches(2.3), Inches(8), Inches(1.1),
             font_size=44, bold=True, color=WHITE)
    add_text(slide, "Intelligence System",
             Inches(0.5), Inches(3.3), Inches(8), Inches(0.9),
             font_size=38, bold=True, color=TEAL)
    add_text(slide, "AI-Powered Forecasting · Full-Stack Dashboard · Machine Learning",
             Inches(0.5), Inches(4.35), Inches(7.8), Inches(0.5),
             font_size=13, color=SLATE_400)

    # Right panel metrics
    metrics = [
        ("Model Accuracy", "99.3% R²", TEAL),
        ("Coverage",       "1990–2031", EMERALD),
        ("Gas Types",      "CO₂ · CH₄ · N₂O", AMBER),
        ("Countries",      "200+", ACCENT_VIOLET),
    ]
    for i, (lbl, val, clr) in enumerate(metrics):
        y = Inches(1.1) + i * Inches(1.5)
        add_text(slide, lbl.upper(), Inches(9.0), y, Inches(4), Inches(0.3),
                 font_size=8, bold=True, color=SLATE_400)
        add_text(slide, val, Inches(9.0), y + Inches(0.3), Inches(4), Inches(0.65),
                 font_size=24, bold=True, color=clr)
        add_section_line(slide, Inches(9.0), y + Inches(1.1), Inches(3.8),
                         RGBColor(0x1E, 0x29, 0x3B))

    # Bottom footer
    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  2026  ·  POWERED BY ML",
             Inches(0.5), Inches(6.9), Inches(8), Inches(0.3),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))

    return slide


# ─── Slide 2 — Problem & Dataset ─────────────────────────────────────────────
def slide_2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, AMBER)

    add_pill_tag(slide, "SLIDE 01  |  PROBLEM & DATA", Inches(0.5), Inches(0.4),
                 RGBColor(0x3D, 0x2C, 0x01), AMBER)
    add_text(slide, "The Problem & Dataset",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    add_section_line(slide, Inches(0.5), Inches(1.6), Inches(12.3), AMBER)

    # Problem card (left)
    add_filled_rect(slide, Inches(0.5), Inches(1.75), Inches(5.9), Inches(4.9), BG_CARD)
    add_text(slide, "🌍  The Challenge",
             Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.5),
             font_size=14, bold=True, color=AMBER)
    bullets = [
        "Greenhouse gas emissions (CO₂, CH₄, N₂O) are the primary driver of climate change.",
        "Policy makers lack precise, country-level forecasts beyond historical records.",
        "Existing tools are complex, not interactive, and don't provide risk classification.",
        "Need: a simple, real-time GHG predictor accessible to everyone.",
    ]
    for i, b in enumerate(bullets):
        add_text(slide, f"▸  {b}",
                 Inches(0.7), Inches(2.55) + i * Inches(0.8), Inches(5.5), Inches(0.7),
                 font_size=11, color=SLATE_300)

    # Dataset card (right)
    add_filled_rect(slide, Inches(6.7), Inches(1.75), Inches(6.1), Inches(4.9), BG_CARD)
    add_text(slide, "📊  Dataset — FAO Global GHG",
             Inches(6.9), Inches(1.9), Inches(5.7), Inches(0.5),
             font_size=14, bold=True, color=TEAL)
    ds_points = [
        ("Source",      "UN FAO (Food & Agriculture Org.)"),
        ("Records",     "~16,000+ rows after cleaning"),
        ("Span",        "1990 – 2021 (30 years)"),
        ("Features",    "Area, Year, Emission Element"),
        ("Targets",     "CO₂, CH₄, N₂O  (kilotonnes)"),
        ("Countries",   "200+ nations covered"),
    ]
    for i, (k, v) in enumerate(ds_points):
        y = Inches(2.55) + i * Inches(0.68)
        add_text(slide, k.upper(), Inches(6.9), y, Inches(1.6), Inches(0.3),
                 font_size=8, bold=True, color=SLATE_400)
        add_text(slide, v, Inches(8.6), y, Inches(4), Inches(0.3),
                 font_size=11, bold=True, color=WHITE)

    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  PROBLEM & DATA",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.25),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))
    return slide


# ─── Slide 3 — EDA & Notebook ─────────────────────────────────────────────────
def slide_3_notebook(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, ACCENT_VIOLET)

    add_pill_tag(slide, "SLIDE 02  |  NOTEBOOK & EDA", Inches(0.5), Inches(0.4),
                 RGBColor(0x2D, 0x1E, 0x52), ACCENT_VIOLET)
    add_text(slide, "Exploratory Data Analysis",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    add_section_line(slide, Inches(0.5), Inches(1.6), Inches(12.3), ACCENT_VIOLET)

    steps = [
        ("01", "Data Loading & Inspection",
         "Loaded FAO CSV into Pandas. Inspected shape, dtypes, and missing values. Dataset had wide format (Year columns) → melted to long format.",
         TEAL),
        ("02", "Cleaning & Preprocessing",
         "Dropped NaN emission rows. Renamed columns. Label-encoded Area and Element categorical columns using LabelEncoder. Applied StandardScaler on features.",
         EMERALD),
        ("03", "EDA — Trends & Distributions",
         "Plotted per-country emission timelines. Aggregated decade averages. Identified top emitters. Visualised CO₂ vs CH₄ vs N₂O distributions with matplotlib/seaborn.",
         AMBER),
        ("04", "Feature Engineering",
         "Built 3 numerical features: [Area_encoded, Element_encoded, Year]. Target: emission value in kilotonnes. Train/test split: 80/20 ratio, random_state=42.",
         ACCENT_VIOLET),
    ]

    cols = [(Inches(0.5), Inches(1.75)), (Inches(6.9), Inches(1.75)),
            (Inches(0.5), Inches(4.3)), (Inches(6.9), Inches(4.3))]

    for i, (num, title, desc, clr) in enumerate(steps):
        lx, ly = cols[i]
        add_filled_rect(slide, lx, ly, Inches(6.0), Inches(2.3), BG_CARD)
        add_text(slide, num, lx + Inches(0.18), ly + Inches(0.12), Inches(0.8), Inches(0.5),
                 font_size=28, bold=True, color=clr)
        add_text(slide, title, lx + Inches(0.18), ly + Inches(0.6), Inches(5.5), Inches(0.4),
                 font_size=13, bold=True, color=WHITE)
        add_text(slide, desc, lx + Inches(0.18), ly + Inches(1.05), Inches(5.6), Inches(1.1),
                 font_size=9.5, color=SLATE_300)

    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  EDA & NOTEBOOK",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.25),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))
    return slide


# ─── Slide 4 — ML Model ───────────────────────────────────────────────────────
def slide_4_model(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, EMERALD)

    add_pill_tag(slide, "SLIDE 03  |  MACHINE LEARNING MODEL", Inches(0.5), Inches(0.4),
                 RGBColor(0x05, 0x3A, 0x25), EMERALD)
    add_text(slide, "Model Training & Performance",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    add_section_line(slide, Inches(0.5), Inches(1.6), Inches(12.3), EMERALD)

    # Left: model details
    add_filled_rect(slide, Inches(0.5), Inches(1.75), Inches(5.9), Inches(5.3), BG_CARD)
    add_text(slide, "🌲  Random Forest Regressor",
             Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.45),
             font_size=15, bold=True, color=EMERALD)

    model_items = [
        ("Algorithm", "RandomForestRegressor (scikit-learn)"),
        ("n_estimators", "100 trees"),
        ("max_depth", "None (fully grown)"),
        ("Features Used", "Area, Gas Type, Year (3 features)"),
        ("Training Split", "80% train / 20% test"),
        ("Data Size", "~16,000 samples after cleaning"),
        ("Export Format", "Pickle (.pkl) for Flask API"),
    ]
    for i, (k, v) in enumerate(model_items):
        y = Inches(2.5) + i * Inches(0.56)
        add_text(slide, f"{k}:", Inches(0.7), y, Inches(1.9), Inches(0.4),
                 font_size=9, bold=True, color=SLATE_400)
        add_text(slide, v, Inches(2.65), y, Inches(3.4), Inches(0.4),
                 font_size=10, color=WHITE)

    # Right: metrics + forecast logic
    add_filled_rect(slide, Inches(6.7), Inches(1.75), Inches(6.1), Inches(2.4), BG_CARD)
    add_text(slide, "📈  Model Performance Metrics",
             Inches(6.9), Inches(1.9), Inches(5.7), Inches(0.45),
             font_size=14, bold=True, color=TEAL)
    kpis = [
        ("R² Score", "99.3%", TEAL),
        ("MAE", "Very Low", EMERALD),
        ("RMSE", "Very Low", EMERALD),
    ]
    for i, (lbl, val, clr) in enumerate(kpis):
        x = Inches(6.9) + i * Inches(2.0)
        add_text(slide, lbl.upper(), x, Inches(2.5), Inches(1.8), Inches(0.3),
                 font_size=8, bold=True, color=SLATE_400)
        add_text(slide, val, x, Inches(2.85), Inches(1.8), Inches(0.5),
                 font_size=18, bold=True, color=clr)

    add_filled_rect(slide, Inches(6.7), Inches(4.35), Inches(6.1), Inches(2.7), BG_CARD)
    add_text(slide, "🔮  Smart Forecast Logic (2022–2031)",
             Inches(6.9), Inches(4.5), Inches(5.7), Inches(0.45),
             font_size=13, bold=True, color=AMBER)
    forecast_desc = (
        "Model was trained on 1990–2021 historical data. For future years:\n\n"
        "▸ RF baseline prediction anchors the starting point.\n"
        "▸ Historical trend slope is calculated from the last 10 years.\n"
        "▸ Forecast = RF Baseline + (Slope × Years Ahead × 0.75)\n"
        "▸ Clipped at 0 to prevent negative emission values."
    )
    add_text(slide, forecast_desc,
             Inches(6.9), Inches(5.05), Inches(5.7), Inches(1.8),
             font_size=9.5, color=SLATE_300)

    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  ML MODEL",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.25),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))
    return slide


# ─── Slide 5 — Backend + Frontend Architecture ───────────────────────────────
def slide_5_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, TEAL)

    add_pill_tag(slide, "SLIDE 04  |  ARCHITECTURE", Inches(0.5), Inches(0.4),
                 RGBColor(0x0D, 0x3D, 0x38), TEAL)
    add_text(slide, "Full-Stack Architecture",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    add_section_line(slide, Inches(0.5), Inches(1.6), Inches(12.3), TEAL)

    # Three columns: Backend, API, Frontend
    cols = [
        ("⚙️  BACKEND", "Python + Flask",
         ["flask, flask-cors, numpy", "pandas, scikit-learn", "pickle (model serialisation)",
          "3 REST endpoints:", "  /predict  →  POST", "  /forecast  →  GET", "  /metadata →  GET"],
         AMBER, Inches(0.5)),
        ("🔗  REST API", "HTTP / JSON",
         ["Area (country name)", "Element (gas type)", "Year (1990–2050)",
          "Responses include:", "  prediction value (kt)", "  insight (Low/Med/High)", "  forecast array"],
         TEAL, Inches(4.95)),
        ("🖥️  FRONTEND", "React + Vite",
         ["react-router-dom navigation", "framer-motion animations", "recharts (Area + Bar charts)",
          "lucide-react icons", "axios for API calls", "Tailwind CSS utility classes",
          "Custom glassmorphism UI"],
         EMERALD, Inches(9.4)),
    ]

    for title, sub, items, clr, lx in cols:
        add_filled_rect(slide, lx, Inches(1.75), Inches(4.1), Inches(5.3), BG_CARD)
        add_text(slide, title, lx + Inches(0.2), Inches(1.9), Inches(3.7), Inches(0.45),
                 font_size=13, bold=True, color=clr)
        add_text(slide, sub, lx + Inches(0.2), Inches(2.4), Inches(3.7), Inches(0.35),
                 font_size=10, color=SLATE_400)
        add_section_line(slide, lx + Inches(0.2), Inches(2.82), Inches(3.7), clr)
        for j, item in enumerate(items):
            add_text(slide, item, lx + Inches(0.2), Inches(2.98) + j * Inches(0.55),
                     Inches(3.7), Inches(0.5), font_size=9.5, color=SLATE_300)

    # Arrow connectors (as thin teal bars)
    for ax in [Inches(4.63), Inches(9.08)]:
        add_filled_rect(slide, ax, Inches(4.4), Inches(0.28), Inches(0.04), TEAL)
        add_text(slide, "→", ax - Inches(0.02), Inches(4.25), Inches(0.35), Inches(0.35),
                 font_size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  ARCHITECTURE",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.25),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))
    return slide


# ─── Slide 6 — UI Design & Final Demo ────────────────────────────────────────
def slide_6_ui_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)
    add_filled_rect(slide, 0, 0, Inches(0.08), SLIDE_H, RED)

    add_pill_tag(slide, "SLIDE 05  |  UI DESIGN & RESULTS", Inches(0.5), Inches(0.4),
                 RGBColor(0x3D, 0x07, 0x07), RED)
    add_text(slide, "UI Design & Key Results",
             Inches(0.5), Inches(0.85), Inches(12), Inches(0.7),
             font_size=32, bold=True, color=WHITE)
    add_section_line(slide, Inches(0.5), Inches(1.6), Inches(12.3), RED)

    # Landing page features (left)
    add_filled_rect(slide, Inches(0.5), Inches(1.75), Inches(5.9), Inches(2.4), BG_CARD)
    add_text(slide, "🏠  Landing Page",
             Inches(0.7), Inches(1.9), Inches(5.5), Inches(0.4),
             font_size=14, bold=True, color=TEAL)
    lp_items = [
        "Dark glassmorphism hero with animated Atmosphere particles",
        "3D card tilt on hover via framer-motion (rotateX/rotateY)",
        "Key metrics strip: Accuracy, Gas Types, Coverage",
        "Gradient CTA button → navigates to Predictor",
    ]
    for i, item in enumerate(lp_items):
        add_text(slide, f"▸  {item}",
                 Inches(0.7), Inches(2.45) + i * Inches(0.42), Inches(5.5), Inches(0.4),
                 font_size=9.5, color=SLATE_300)

    # Dashboard features (right)
    add_filled_rect(slide, Inches(6.7), Inches(1.75), Inches(6.1), Inches(2.4), BG_CARD)
    add_text(slide, "📊  GHG Dashboard (Predictor)",
             Inches(6.9), Inches(1.9), Inches(5.7), Inches(0.4),
             font_size=14, bold=True, color=AMBER)
    db_items = [
        "KPI row: Peak, Avg Historical, 2031 Trend, R² Score",
        "Area chart: historical timeline + AI forecast overlay",
        "Decade bar chart: 1990s → 2020s averages",
        "Threat level badge: 🟢 Low · 🟡 Medium · 🔴 High",
        "Mitigation strategy panel specific to gas type",
    ]
    for i, item in enumerate(db_items):
        add_text(slide, f"▸  {item}",
                 Inches(6.9), Inches(2.45) + i * Inches(0.42), Inches(5.7), Inches(0.4),
                 font_size=9.5, color=SLATE_300)

    # Results / Outcomes strip
    add_filled_rect(slide, Inches(0.5), Inches(4.35), Inches(12.3), Inches(2.7), BG_CARD)
    add_text(slide, "🏆  Key Outcomes & Achievements",
             Inches(0.7), Inches(4.5), Inches(11), Inches(0.45),
             font_size=14, bold=True, color=EMERALD)
    add_section_line(slide, Inches(0.7), Inches(4.98), Inches(11.8), EMERALD)

    outcomes = [
        ("99.3% R²", "Model Accuracy\non test set", TEAL),
        ("3 Endpoints", "REST API routes:\npredict · forecast · metadata", AMBER),
        ("5 UI Screens", "Landing + Dashboard\nwith full interactivity", EMERALD),
        ("30 yrs Data", "Historical FAO dataset\n1990 – 2021", ACCENT_VIOLET),
        ("10 yr Forecast", "AI-powered projection\n2022 – 2031", RED),
    ]
    for i, (val, desc, clr) in enumerate(outcomes):
        x = Inches(0.9) + i * Inches(2.38)
        add_text(slide, val, x, Inches(5.1), Inches(2.2), Inches(0.55),
                 font_size=20, bold=True, color=clr)
        add_text(slide, desc, x, Inches(5.65), Inches(2.2), Inches(0.7),
                 font_size=8.5, color=SLATE_400)

    add_text(slide, "GLOBAL GHG INTELLIGENCE  ·  UI & RESULTS",
             Inches(0.5), Inches(7.0), Inches(12), Inches(0.25),
             font_size=7, color=RGBColor(0x33, 0x4D, 0x5C))
    return slide


# ─── Slide 7 (Bonus / Conclusion) — Summary & Thank You ──────────────────────
def slide_7_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, BG_DARK)

    # Full-bleed teal top bar
    add_filled_rect(slide, 0, 0, SLIDE_W, Inches(0.18), TEAL)
    add_filled_rect(slide, 0, Inches(7.33), SLIDE_W, Inches(0.18), TEAL)

    add_text(slide, "Thank You",
             0, Inches(1.2), SLIDE_W, Inches(1.2),
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Global GHG Emissions Intelligence System",
             0, Inches(2.5), SLIDE_W, Inches(0.7),
             font_size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "A full-stack AI project: from raw FAO data to interactive dashboard",
             0, Inches(3.2), SLIDE_W, Inches(0.5),
             font_size=13, color=SLATE_400, align=PP_ALIGN.CENTER, italic=True)

    # Summary pills
    tags = [
        ("Python / Pandas", AMBER),
        ("scikit-learn RF", EMERALD),
        ("Flask REST API", TEAL),
        ("React + Vite", ACCENT_VIOLET),
        ("framer-motion", RED),
        ("recharts", AMBER),
    ]
    total_w = Inches(11)
    tag_w = Inches(1.75)
    gap = Inches(0.25)
    start_x = (SLIDE_W - (len(tags) * tag_w + (len(tags) - 1) * gap)) / 2

    for i, (tag, clr) in enumerate(tags):
        x = start_x + i * (tag_w + gap)
        shape = slide.shapes.add_shape(5, x, Inches(4.2), tag_w, Inches(0.38))
        shape.adjustments[0] = 0.5
        faded_bg = RGBColor(
            min(255, clr[0] // 5),
            min(255, clr[1] // 5),
            min(255, clr[2] // 5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = faded_bg
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = tag
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = clr
        run.font.name = "Calibri"

    add_text(slide, "BUILT FOR CAPSTONE DEMONSTRATION  ·  2026",
             0, Inches(6.8), SLIDE_W, Inches(0.3),
             font_size=8, color=RGBColor(0x33, 0x4D, 0x5C), align=PP_ALIGN.CENTER)

    return slide


# ─── Main ─────────────────────────────────────────────────────────────────────
def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_notebook(prs)
    slide_4_model(prs)
    slide_5_architecture(prs)
    slide_6_ui_demo(prs)
    slide_7_conclusion(prs)

    out_path = r"c:\Users\HP\Downloads\Capstone Project second\GHG_Capstone_Presentation.pptx"
    prs.save(out_path)
    print(f"\n✅  Presentation saved to:\n    {out_path}\n")
    print("📊  Slides generated:")
    print("    1. Title / Hero")
    print("    2. Problem & Dataset")
    print("    3. EDA & Notebook")
    print("    4. ML Model & Training")
    print("    5. Full-Stack Architecture")
    print("    6. UI Design & Results")
    print("    7. Thank You / Conclusion")


if __name__ == "__main__":
    build_presentation()
